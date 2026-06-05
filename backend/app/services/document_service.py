"""
Mela AI - Document Processing Service

Local text extraction from documents (PDF, Word, Excel, PPT, CSV, JSON, HTML,
code files, etc.). No Azure Storage dependency — blob upload is handled
separately and only when configured.
"""

import csv
import hashlib
import io
import json
import logging
import re
from pathlib import Path
from typing import Any, Dict, Optional, Tuple

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Local document text extractor. No cloud dependencies."""

    # Maps MIME type → short type tag
    MIME_TO_TYPE: Dict[str, str] = {
        "application/pdf": "pdf",
        "application/vnd.openxmlformats-officedocument"
        ".wordprocessingml.document": "docx",
        "application/msword": "docx",
        "application/vnd.openxmlformats-officedocument"
        ".spreadsheetml.sheet": "xlsx",
        "application/vnd.ms-excel": "xlsx",
        "application/vnd.openxmlformats-officedocument"
        ".presentationml.presentation": "pptx",
        "application/vnd.ms-powerpoint": "pptx",
        # OpenDocument formats (LibreOffice / OpenOffice)
        "application/vnd.oasis.opendocument.text": "odt",
        "application/vnd.oasis.opendocument.spreadsheet": "ods",
        "application/vnd.oasis.opendocument.presentation": "odp",
        # EPUB
        "application/epub+zip": "epub",
        "text/plain": "txt",
        "text/csv": "csv",
        "text/markdown": "md",
        "text/html": "html",
        "application/json": "json",
        "application/xml": "xml",
        "text/xml": "xml",
        "application/rtf": "rtf",
        "text/rtf": "rtf",
        "application/x-rtf": "rtf",
        "message/rfc822": "eml",
        "application/zip": "zip",
        "application/x-zip-compressed": "zip",
        "application/x-zip": "zip",
        "text/javascript": "code",
        "application/javascript": "code",
        "text/x-python": "code",
        "application/x-python-code": "code",
        "application/x-sh": "code",
        "text/x-shellscript": "code",
        "application/x-yaml": "code",
        "text/yaml": "code",
        "text/x-yaml": "code",
        "application/x-toml": "code",
        "text/x-sql": "code",
        "application/sql": "code",
    }

    # Extension fallbacks (when MIME type is generic or missing)
    EXT_TO_TYPE: Dict[str, str] = {
        "pdf": "pdf",
        "docx": "docx", "doc": "docx",
        "xlsx": "xlsx", "xls": "xlsx",
        "pptx": "pptx", "ppt": "pptx",
        # OpenDocument
        "odt": "odt", "ott": "odt",
        "ods": "ods", "ots": "ods",
        "odp": "odp", "otp": "odp",
        # EPUB
        "epub": "epub",
        "txt": "txt",
        "md": "md", "markdown": "md",
        "csv": "csv",
        "json": "json",
        "xml": "xml",
        "html": "html", "htm": "html",
        # Code / config
        "py": "code", "js": "code", "ts": "code",
        "jsx": "code", "tsx": "code",
        "java": "code", "cpp": "code", "c": "code",
        "cs": "code", "go": "code", "rs": "code",
        "rb": "code", "php": "code",
        "sh": "code", "bash": "code", "zsh": "code",
        "yaml": "code", "yml": "code",
        "toml": "code", "ini": "code", "cfg": "code",
        "sql": "code",
        "r": "code", "jl": "code",
        "rtf": "rtf",
        "eml": "eml", "email": "eml",
        "zip": "zip",
        "swift": "code", "kt": "code", "kts": "code",
        "dart": "code", "scala": "code",
        "log": "txt", "conf": "code", "env": "code",
        "tex": "code", "rst": "md",
    }

    def detect_type(self, content_type: str = "", filename: str = "") -> str:
        """Map MIME type / filename extension → short type tag."""
        # Strip MIME parameters (e.g. "text/plain; charset=utf-8")
        ct = (content_type or "").split(";")[0].strip().lower()
        if ct in self.MIME_TO_TYPE:
            return self.MIME_TO_TYPE[ct]
        # text/* catch-all
        if ct.startswith("text/"):
            return "txt"
        # Extension fallback
        ext = Path(filename).suffix.lower().lstrip(".")
        return self.EXT_TO_TYPE.get(ext, "unknown")

    def extract_text(
        self,
        file_data: bytes,
        content_type: str,
        filename: str = "",
    ) -> Tuple[str, Dict[str, Any]]:
        """
        Extract text from a file.

        Args:
            file_data:    Raw file bytes.
            content_type: Full MIME type string (e.g. "application/pdf").
            filename:     Original filename (used as extension fallback).

        Returns:
            (extracted_text, metadata_dict)
        """
        file_type = self.detect_type(content_type, filename)
        metadata: Dict[str, Any] = {
            "filename": filename,
            "content_type": content_type,
            "file_type": file_type,
            "file_size": len(file_data),
        }

        try:
            if file_type == "pdf":
                text, meta = self._extract_pdf(file_data)
            elif file_type == "docx":
                text, meta = self._extract_docx(file_data)
            elif file_type == "xlsx":
                text, meta = self._extract_xlsx(file_data)
            elif file_type == "pptx":
                text, meta = self._extract_pptx(file_data)
            elif file_type == "csv":
                text, meta = self._extract_csv(file_data)
            elif file_type == "json":
                text, meta = self._extract_json(file_data)
            elif file_type == "html":
                text, meta = self._extract_html(file_data)
            elif file_type == "rtf":
                text, meta = self._extract_rtf(file_data)
            elif file_type == "eml":
                text, meta = self._extract_eml(file_data)
            elif file_type == "zip":
                text, meta = self._extract_zip(file_data)
            elif file_type == "odt":
                text, meta = self._extract_odf(file_data, "content.xml")
            elif file_type in ("ods", "odp"):
                text, meta = self._extract_odf(file_data, "content.xml")
            elif file_type == "epub":
                text, meta = self._extract_epub(file_data)
            elif file_type in ("txt", "md", "code", "xml"):
                text = file_data.decode("utf-8", errors="replace")
                meta = {}
            else:
                logger.warning(
                    "Unsupported type %r (%r) — attempting best-effort extraction",
                    file_type, content_type,
                )
                text, meta = self._extract_unknown(file_data, filename)

            metadata.update(meta)
            return text, metadata

        except Exception as exc:
            logger.error(f"Text extraction error [{filename!r}]: {exc}")
            return "", metadata

    # ── Extractors ────────────────────────────────────────────────────────────

    def _extract_pdf(self, file_data: bytes) -> Tuple[str, Dict]:
        import fitz  # PyMuPDF

        doc = fitz.open(stream=file_data, filetype="pdf")
        parts = []
        for page_num, page in enumerate(doc):
            t = page.get_text()
            if t.strip():
                parts.append(f"[Page {page_num + 1}]\n{t}")
        meta = {
            "page_count": len(doc),
            "title": doc.metadata.get("title", ""),
            "author": doc.metadata.get("author", ""),
        }
        doc.close()
        return "\n\n".join(parts), meta

    def _extract_docx(self, file_data: bytes) -> Tuple[str, Dict]:
        from docx import Document as DocxDocument

        doc = DocxDocument(io.BytesIO(file_data))
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text for cell in row.cells)
                if row_text.strip():
                    parts.append(row_text)
        meta = {
            "paragraph_count": len(doc.paragraphs),
            "table_count": len(doc.tables),
        }
        return "\n\n".join(parts), meta

    def _extract_xlsx(self, file_data: bytes) -> Tuple[str, Dict]:
        from openpyxl import load_workbook

        MAX_ROWS_PER_SHEET = 5000  # Prevent hang on huge spreadsheets

        wb = load_workbook(io.BytesIO(file_data), data_only=True)
        parts = []
        total_rows = 0
        truncated_sheets = []
        hidden_sheets = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            # Track hidden sheets
            if sheet.sheet_state != "visible":
                hidden_sheets.append(sheet_name)
                parts.append(f"[Sheet: {sheet_name} (hidden)]")
            else:
                parts.append(f"[Sheet: {sheet_name}]")
            row_count = 0
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join(
                    str(c) if c is not None else "" for c in row
                )
                if row_text.strip():
                    parts.append(row_text)
                    row_count += 1
                    total_rows += 1
                if row_count >= MAX_ROWS_PER_SHEET:
                    truncated_sheets.append(sheet_name)
                    parts.append(f"[... truncated at {MAX_ROWS_PER_SHEET} rows]")
                    break

        # Check for formulas with missing values (data_only=True returns None
        # for formulas that were never calculated)
        has_formulas = False
        try:
            wb_formulas = load_workbook(io.BytesIO(file_data), data_only=False)
            for sn in wb_formulas.sheetnames:
                fs = wb_formulas[sn]
                for row in fs.iter_rows(min_row=1, max_row=5, values_only=True):
                    for cell in row:
                        if isinstance(cell, str) and cell.startswith("="):
                            has_formulas = True
                            break
                    if has_formulas:
                        break
                if has_formulas:
                    break
        except Exception:
            pass

        meta = {
            "sheet_count": len(wb.sheetnames),
            "sheets": list(wb.sheetnames),
            "total_rows": total_rows,
            "has_formulas": has_formulas,
        }
        if truncated_sheets:
            meta["truncated_sheets"] = truncated_sheets
        if hidden_sheets:
            meta["hidden_sheets"] = hidden_sheets

        # ── Build DATA CARDs per sheet via pandas ─────────────────────────
        # Note: we build DataFrames from openpyxl row iteration rather than
        # calling pd.read_excel() because the latter pins a newer openpyxl
        # version than we currently ship.
        try:
            import pandas as pd  # type: ignore
            cards: list = []
            profiles: list = []
            for sn in wb.sheetnames[:10]:  # cap at 10 sheets
                try:
                    sh = wb[sn]
                    if sh.sheet_state != "visible":
                        continue
                    # Use sh.max_row so we don't fabricate thousands of empty
                    # rows on sparse sheets.
                    row_limit = min(sh.max_row or 0, 10000)
                    if row_limit <= 0:
                        continue
                    raw = list(sh.iter_rows(values_only=True, max_row=row_limit))
                    # Drop rows that are entirely empty
                    raw = [r for r in raw
                           if any(v is not None and v != "" for v in r)]
                    if not raw:
                        continue
                    header = [
                        str(h) if h is not None else f"col_{i}"
                        for i, h in enumerate(raw[0])
                    ]
                    body = raw[1:] if len(raw) > 1 else []
                    df = pd.DataFrame(body, columns=header)
                    # Promote numeric-looking object columns. Manual try/except
                    # because pandas removed errors='ignore' in 3.x.
                    for col in df.columns:
                        if df[col].dtype == object:
                            try:
                                converted = pd.to_numeric(df[col], errors="raise")
                                df[col] = converted
                            except (ValueError, TypeError):
                                pass
                except Exception:
                    continue
                if df.empty:
                    continue
                card, profile = self._dataframe_card_from_df(df, sheet_label=sn)
                if card:
                    cards.append(card)
                if profile:
                    profiles.append(profile)
            if cards:
                parts = cards + [""] + parts
            if profiles:
                meta["structured_profile"] = profiles
        except Exception as exc:
            logger.debug("xlsx DATA CARD build failed: %s", exc)

        return "\n".join(parts), meta

    def _extract_pptx(self, file_data: bytes) -> Tuple[str, Dict]:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(file_data))
        parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            texts = [
                shape.text
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            if texts:
                parts.append(f"[Slide {slide_num}]\n" + "\n".join(texts))
        meta = {"slide_count": len(prs.slides)}
        return "\n\n".join(parts), meta

    def _extract_csv(self, file_data: bytes) -> Tuple[str, Dict]:
        """Extract CSV + a pandas-derived DATA CARD (shape, dtypes, head, stats).

        The DATA CARD makes the file immediately queryable by the LLM and
        improves embedding quality because column names and numeric
        distributions are present in-band with the raw rows.
        """
        content = file_data.decode("utf-8", errors="replace")
        # Raw-rows fallback (works even if pandas fails on weird CSV dialects)
        reader = csv.reader(io.StringIO(content))
        rows = list(reader)
        raw_text = "\n".join(" | ".join(r) for r in rows)

        data_card, profile = self._build_dataframe_card(content, sheet_label="CSV")
        full = (data_card + "\n\n" + raw_text) if data_card else raw_text
        meta: Dict[str, Any] = {"row_count": len(rows)}
        if profile:
            meta["structured_profile"] = profile
        return full, meta

    def _build_dataframe_card(
        self, csv_text: str, sheet_label: str = "Data"
    ) -> Tuple[str, Dict[str, Any]]:
        """Return (markdown DATA CARD, machine-readable profile dict)."""
        try:
            import pandas as pd  # type: ignore
        except Exception:
            return "", {}
        try:
            df = pd.read_csv(io.StringIO(csv_text), engine="python",
                             on_bad_lines="skip", nrows=10000)
        except Exception as exc:
            logger.debug("pandas CSV parse failed for %s: %s", sheet_label, exc)
            return "", {}
        return self._dataframe_card_from_df(df, sheet_label)

    def _dataframe_card_from_df(
        self, df, sheet_label: str = "Data"
    ) -> Tuple[str, Dict[str, Any]]:
        try:
            import pandas as pd  # type: ignore  # noqa: F401
        except Exception:
            return "", {}
        rows, cols = df.shape
        col_profile: list = []
        for c in df.columns[:50]:  # cap wide sheets
            series = df[c]
            dtype = str(series.dtype)
            non_null = int(series.notna().sum())
            entry = {"name": str(c), "dtype": dtype, "non_null": non_null}
            try:
                if series.dtype.kind in "biufc":  # numeric
                    entry["min"] = float(series.min()) if non_null else None
                    entry["max"] = float(series.max()) if non_null else None
                    entry["mean"] = float(series.mean()) if non_null else None
                else:
                    uniq = series.dropna().astype(str).unique()[:10]
                    entry["sample_values"] = [str(v)[:60] for v in uniq]
                    entry["unique_count"] = int(series.nunique(dropna=True))
            except Exception:
                pass
            col_profile.append(entry)

        head_md = ""
        try:
            head_md = df.head(10).to_markdown(index=False)
        except Exception:
            try:
                head_md = df.head(10).to_string(index=False)
            except Exception:
                head_md = ""

        lines = [
            f"[DATA CARD — {sheet_label}]",
            f"Shape: {rows} rows × {cols} columns",
            "Columns:",
        ]
        for c in col_profile:
            extras = []
            if "min" in c and c["min"] is not None:
                extras.append(
                    f"range={c['min']:.4g}..{c['max']:.4g} "
                    f"mean={c['mean']:.4g}"
                )
            elif "sample_values" in c:
                extras.append(f"uniq={c['unique_count']} "
                              f"samples={c['sample_values'][:3]}")
            lines.append(f"  - {c['name']} ({c['dtype']}) "
                         f"non-null={c['non_null']} "
                         + " ".join(extras))
        if head_md:
            lines.append("")
            lines.append("First rows:")
            lines.append(head_md)
        card = "\n".join(lines)
        profile = {
            "sheet": sheet_label,
            "shape": [rows, cols],
            "columns": col_profile,
        }
        return card, profile

    def _extract_json(self, file_data: bytes) -> Tuple[str, Dict]:
        try:
            data = json.loads(file_data.decode("utf-8", errors="replace"))
            text = json.dumps(data, indent=2, ensure_ascii=False)
        except Exception:
            text = file_data.decode("utf-8", errors="replace")
        return text, {}

    def _extract_html(self, file_data: bytes) -> Tuple[str, Dict]:
        """Extract visible text from HTML using the stdlib html.parser."""
        from html.parser import HTMLParser

        raw = file_data.decode("utf-8", errors="replace")

        class _TextCollector(HTMLParser):
            _SKIP = {"script", "style", "head", "meta", "link", "noscript"}

            def __init__(self):
                super().__init__()
                self._parts: list = []
                self._skip_depth: int = 0

            def handle_starttag(self, tag, attrs):
                if tag.lower() in self._SKIP:
                    self._skip_depth += 1

            def handle_endtag(self, tag):
                if tag.lower() in self._SKIP and self._skip_depth > 0:
                    self._skip_depth -= 1

            def handle_data(self, data):
                if self._skip_depth == 0:
                    stripped = data.strip()
                    if stripped:
                        self._parts.append(stripped)

        try:
            parser = _TextCollector()
            parser.feed(raw)
            text = " ".join(parser._parts)
            text = re.sub(r"\s+", " ", text).strip()
        except Exception:
            text = re.sub(r"<[^>]+>", " ", raw)
            text = re.sub(r"\s+", " ", text).strip()
        return text, {}

    def _extract_rtf(self, file_data: bytes) -> Tuple[str, Dict]:
        """Strip RTF control words and extract plain text."""
        raw = file_data.decode("latin-1", errors="replace")
        try:
            # Remove RTF control words and groups
            text = re.sub(r"\\[a-z]+\d* ?", " ", raw)
            text = re.sub(r"[{}\\]", " ", text)
            # Decode RTF hex escapes like \'e9
            text = re.sub(
                r"\\'([0-9a-fA-F]{2})",
                lambda m: chr(int(m.group(1), 16)),
                text,
            )
            text = re.sub(r"\s+", " ", text).strip()
        except Exception:
            text = raw
        return text, {"format": "rtf"}

    def _extract_eml(self, file_data: bytes) -> Tuple[str, Dict]:
        """Extract text and metadata from .eml email files."""
        import email as email_lib
        msg = email_lib.message_from_bytes(file_data)
        subject = msg.get("Subject", "")
        from_ = msg.get("From", "")
        to_ = msg.get("To", "")
        date_ = msg.get("Date", "")
        parts = [
            f"Subject: {subject}",
            f"From: {from_}",
            f"To: {to_}",
            f"Date: {date_}",
            "---",
        ]
        for part in msg.walk():
            ct = part.get_content_type()
            if ct == "text/plain":
                payload = part.get_payload(decode=True)
                if payload:
                    parts.append(payload.decode("utf-8", errors="replace"))
                    break
            elif ct == "text/html" and len(parts) <= 6:
                payload = part.get_payload(decode=True)
                if payload:
                    raw_html = payload.decode("utf-8", errors="replace")
                    plain = re.sub(r"<[^>]+>", " ", raw_html)
                    plain = re.sub(r"\s+", " ", plain).strip()
                    parts.append(plain)
        return "\n".join(parts), {
            "subject": subject, "from": from_, "to": to_,
        }

    def _extract_zip(self, file_data: bytes) -> Tuple[str, Dict]:
        """Extract text from readable files inside a ZIP archive."""
        import zipfile

        text_parts: list = []
        file_count = 0
        text_ext = {
            ".txt", ".md", ".csv", ".json", ".xml",
            ".html", ".htm", ".py", ".js", ".ts",
            ".sql", ".yaml", ".yml", ".toml", ".ini",
            ".cfg", ".sh", ".rb", ".go", ".rs",
        }
        try:
            with zipfile.ZipFile(io.BytesIO(file_data)) as zf:
                for name in zf.namelist()[:50]:  # cap at 50 entries
                    ext = Path(name).suffix.lower()
                    if ext in text_ext:
                        try:
                            raw = zf.read(name)
                            content = raw.decode("utf-8", errors="replace")
                            text_parts.append(f"=== {name} ===\n{content}")
                            file_count += 1
                        except Exception:
                            pass
        except Exception as e:
            return f"Could not open ZIP: {e}", {}
        if text_parts:
            text = "\n\n".join(text_parts)
        else:
            text = "(No readable text files found in archive)"
        return text, {"zip_files_extracted": file_count}

    def _extract_odf(
        self, file_data: bytes, content_file: str = "content.xml"
    ) -> Tuple[str, Dict]:
        """
        Extract plain text from OpenDocument Format files (ODT/ODS/ODP).
        These are ZIP archives containing an XML content file.
        """
        import zipfile
        from xml.etree import ElementTree

        try:
            with zipfile.ZipFile(io.BytesIO(file_data)) as zf:
                if content_file not in zf.namelist():
                    return "(Could not read ODF content)", {}
                xml_data = zf.read(content_file)

            root = ElementTree.fromstring(xml_data)
            texts: list = []
            for elem in root.iter():
                if elem.text and elem.text.strip():
                    texts.append(elem.text.strip())
                if elem.tail and elem.tail.strip():
                    texts.append(elem.tail.strip())

            text = "\n".join(texts)
            return text, {"format": "odf"}
        except Exception as e:
            logger.warning("ODF extraction failed: %s", e)
            return "", {}

    def _extract_epub(self, file_data: bytes) -> Tuple[str, Dict]:
        """
        Extract text from EPUB files (ZIP of XHTML/HTML content files).
        Reads all .xhtml/.html items from the EPUB spine.
        """
        import zipfile
        from xml.etree import ElementTree

        parts: list = []
        chapter_count = 0
        try:
            with zipfile.ZipFile(io.BytesIO(file_data)) as zf:
                names = zf.namelist()
                # Collect HTML/XHTML content files (skip nav/toc)
                content_files = sorted(
                    n for n in names
                    if n.endswith((".xhtml", ".html", ".htm"))
                    and "nav" not in n.lower()
                    and "toc" not in n.lower()
                )
                for name in content_files[:100]:  # cap chapters
                    try:
                        raw = zf.read(name).decode("utf-8", errors="replace")
                        # Strip tags
                        text = re.sub(r"<[^>]+>", " ", raw)
                        text = re.sub(r"\s+", " ", text).strip()
                        if text:
                            parts.append(text)
                            chapter_count += 1
                    except Exception:
                        pass
        except Exception as e:
            return f"Could not read EPUB: {e}", {}

        text = "\n\n".join(parts) if parts else "(No readable content in EPUB)"
        return text, {"chapter_count": chapter_count}

    # ── Utilities ─────────────────────────────────────────────────────────────

    def _extract_unknown(
        self, file_data: bytes, filename: str = ""
    ) -> Tuple[str, Dict]:
        """Best-effort extraction for unknown types.

        Strategy:
        1. Detect the magic bytes for common formats and re-route.
        2. Fall back to UTF-8 / latin-1 decode with low-entropy filtering
           so binary garbage does not pollute the index.
        """
        if not file_data:
            return "", {"warning": "empty file"}

        # Magic-byte re-route
        head = file_data[:8]
        if head.startswith(b"%PDF-"):
            try:
                return self._extract_pdf(file_data)
            except Exception:
                pass
        if head.startswith(b"PK\x03\x04"):
            # ZIP-based; try docx / xlsx / pptx / zip in order
            for fn in (self._extract_docx, self._extract_xlsx,
                       self._extract_pptx, self._extract_zip):
                try:
                    t, m = fn(file_data)
                    if t and t.strip():
                        return t, m
                except Exception:
                    continue
        if head[:2] in (b"{\n", b"{\r", b"{ ", b"[{") or head[:1] in (b"{", b"["):
            try:
                return self._extract_json(file_data)
            except Exception:
                pass
        # Text heuristic: > 70% printable ASCII → keep; else mark unreadable.
        try:
            sample = file_data[:4096]
            printable = sum(1 for b in sample if 9 <= b <= 126 or b in (10, 13))
            if sample and printable / len(sample) >= 0.70:
                text = file_data.decode("utf-8", errors="replace")
                return text, {"warning": "extracted as plain text heuristic"}
        except Exception:
            pass
        return (
            f"(Binary file — not text-extractable: {filename or 'file'})",
            {"warning": "binary, skipped", "size": len(file_data)},
        )

    def get_content_hash(self, file_data: bytes) -> str:
        return hashlib.sha256(file_data).hexdigest()

    def detect_file_type(self, filename: str, content_type: str = "") -> str:
        """Legacy helper — prefer detect_type()."""
        return self.detect_type(content_type, filename)


# ── Lazy singleton ─────────────────────────────────────────────────────────────

_document_processor: Optional[DocumentProcessor] = None


def get_document_processor() -> DocumentProcessor:
    """Return the shared DocumentProcessor instance (lazy singleton)."""
    global _document_processor
    if _document_processor is None:
        _document_processor = DocumentProcessor()
    return _document_processor


# Backward-compat alias — callers should use get_document_processor().
document_processor = None


def extract_text(
    file_data: bytes, content_type: str, filename: str = ""
) -> str:
    """Module-level convenience wrapper — returns text only (no metadata)."""
    text, _ = get_document_processor().extract_text(
        file_data, content_type, filename
    )
    return text
